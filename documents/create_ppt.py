# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()

# 设置默认字体
font_name = "Microsoft YaHei"
font_size = Pt(12)

# 定义颜色方案
DARK_BROWN = RGBColor(0x4E, 0x34, 0x2E)
LIGHT_BROWN = RGBColor(0xD7, 0xCC, 0xC8)
OFF_WHITE = RGBColor(0xFF, 0xFC, 0xF9)
GOLD = RGBColor(0xFF, 0xD5, 0x4F)
GREEN = RGBColor(0x8B, 0xC3, 0x4A)

# 封面页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 设置标题
title.text = "咖啡豆详解：从产地到风味的全面指南"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(44)
title_para.font.color.rgb = DARK_BROWN

# 设置副标题
subtitle.text = "探索不同咖啡豆的独特风味"
subtitle_para = subtitle.text_frame.paragraphs[0]
subtitle_para.font.name = "Microsoft YaHei"
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = DARK_BROWN

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 目录页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "目录"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "1. 阿拉比卡 (Arabica)\n2. 罗布斯塔 (Robusta)\n3. 利比里卡 (Liberica)\n4. 埃克塞尔萨 (Excelsa)\n5. 结语"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(24)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 阿拉比卡 (Arabica)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "阿拉比卡 (Arabica)"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "产地：\n- 拉丁美洲（巴西、哥伦比亚）\n- 东非（埃塞俄比亚、肯尼亚）\n- 亚洲部分地区（印度尼西亚）\n\n品种特点：\n- 占全球咖啡产量的60%-70%\n- 豆形较长，含糖量高，咖啡因含量低\n\n风味特点：\n- 口感柔和，酸甜平衡\n- 埃塞俄比亚：浆果和柑橘香气\n- 哥伦比亚：坚果和巧克力味道\n- 巴西：坚果和焦糖风味"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(18)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 罗布斯塔 (Robusta)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "罗布斯塔 (Robusta)"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "产地：\n- 非洲（乌干达、科特迪瓦）\n- 印度尼西亚和越南\n\n品种特点：\n- 占全球咖啡产量的30%-40%\n- 豆形较圆，咖啡因含量高\n\n风味特点：\n- 口感浓烈，苦味明显\n- 木质、泥土和坚果风味"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(18)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 利比里卡 (Liberica)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "利比里卡 (Liberica)"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "产地：\n- 西非（利比里亚、几内亚）\n- 东南亚（菲律宾、马来西亚）\n\n品种特点：\n- 形状不规则，产量稀少\n\n风味特点：\n- 烟熏和木质风味\n- 酸度低，苦味较重"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(18)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 埃克塞尔萨 (Excelsa)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "埃克塞尔萨 (Excelsa)"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "产地：\n- 东南亚（越南、泰国）\n- 非洲部分地区\n\n品种特点：\n- 形状细长，适应性强\n\n风味特点：\n- 风味复杂多变\n- 葡萄、樱桃香气"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(18)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 结语
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题
title.text = "结语"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(36)
title_para.font.color.rgb = DARK_BROWN

# 设置内容
content.text = "不同咖啡豆品种和产地赋予咖啡独特风味\n了解这些信息有助于选择适合自己口味的咖啡\n鼓励尝试不同风味的咖啡"
content_para = content.text_frame.paragraphs[0]
content_para.font.name = "Microsoft YaHei"
content_para.font.size = Pt(24)
content_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 结束页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title

# 设置标题
title.text = "感谢观看"
title_para = title.text_frame.paragraphs[0]
title_para.font.name = "Microsoft YaHei"
title_para.font.size = Pt(44)
title_para.font.color.rgb = DARK_BROWN

# 设置背景颜色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = OFF_WHITE

# 保存演示文稿
prs.save("咖啡豆详解_优化版.pptx")